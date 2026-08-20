from __future__ import annotations

import csv
import importlib.util
import json
from dataclasses import dataclass, field
from pathlib import Path

from .diagram_engine import DiagramEngine, DiagramSpec
from .presentation_builder import PresentationBuilder, SlideDeckSpec
from .schemas import ArtifactSpec
from .spreadsheet_builder import SpreadsheetSpec

DOCX_AVAILABLE = importlib.util.find_spec("docx") is not None
OPENPYXL_AVAILABLE = importlib.util.find_spec("openpyxl") is not None
PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None


@dataclass
class ValidationReport:
    valid: bool
    errors: list[str] = field(default_factory=list)


class ArtifactValidator:
    """File creation is never automatically treated as success. Every
    artifact type gets a real, type-specific check."""

    def validate_markdown(self, spec: ArtifactSpec, path: Path) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        text = path.read_text(encoding="utf-8")
        errors = []
        for section in spec.content_sections:
            if f"## {section.heading}" not in text:
                errors.append(f"Missing required section: {section.heading}")
        if spec.data_sources and "## Sources" not in text:
            errors.append("Sources section missing despite configured data_sources")
        return ValidationReport(not errors, errors)

    def validate_json(self, path: Path) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        try:
            json.loads(path.read_text(encoding="utf-8"))
        except Exception as error:
            return ValidationReport(False, [f"Invalid JSON: {error}"])
        return ValidationReport(True, [])

    def validate_csv(self, path: Path, expected_headers: list[str] | None = None) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        with path.open(encoding="utf-8") as handle:
            rows = list(csv.reader(handle))
        if not rows:
            return ValidationReport(False, ["CSV file is empty"])
        if expected_headers is not None and rows[0] != expected_headers:
            return ValidationReport(False, [f"CSV headers do not match: {rows[0]} != {expected_headers}"])
        return ValidationReport(True, [])

    def validate_diagram(self, diagram_spec: DiagramSpec, path: Path) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        errors = DiagramEngine.validate(diagram_spec, path.read_text(encoding="utf-8"))
        return ValidationReport(not errors, errors)

    def validate_docx(self, path: Path) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        if not DOCX_AVAILABLE:
            return ValidationReport(False, ["python-docx is not installed; cannot verify DOCX opens"])
        try:
            from docx import Document as DocxDocument

            document = DocxDocument(str(path))
            if not document.paragraphs:
                return ValidationReport(False, ["DOCX has no content"])
        except Exception as error:
            return ValidationReport(False, [f"DOCX failed to open: {error}"])
        return ValidationReport(True, [])

    def validate_spreadsheet(self, path: Path, spreadsheet: SpreadsheetSpec) -> ValidationReport:
        if not path.exists():
            return ValidationReport(False, ["File does not exist"])
        if not OPENPYXL_AVAILABLE:
            return ValidationReport(False, ["openpyxl is not installed; cannot verify workbook opens"])
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(str(path))
        except Exception as error:
            return ValidationReport(False, [f"Workbook failed to open: {error}"])
        expected_sheet_names = {sheet.name[:31] for sheet in spreadsheet.sheets}
        missing = expected_sheet_names - set(workbook.sheetnames)
        errors = [f"Missing expected sheet(s): {sorted(missing)}"] if missing else []
        return ValidationReport(not errors, errors)

    def validate_presentation(self, path: Path, deck: SlideDeckSpec) -> ValidationReport:
        structural_errors = PresentationBuilder.validate(deck)
        if not path.exists():
            return ValidationReport(False, ["File does not exist", *structural_errors])
        if not PPTX_AVAILABLE:
            return ValidationReport(False, ["python-pptx is not installed; cannot verify file opens", *structural_errors])
        try:
            from pptx import Presentation

            presentation = Presentation(str(path))
        except Exception as error:
            return ValidationReport(False, [f"Presentation failed to open: {error}", *structural_errors])
        errors = list(structural_errors)
        if len(presentation.slides) != len(deck.slides) + 1:
            errors.append("Rendered slide count does not match the planned deck (title slide + content slides)")
        return ValidationReport(not errors, errors)
