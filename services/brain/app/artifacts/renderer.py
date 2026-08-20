from __future__ import annotations

import csv
import importlib.util
import json
from pathlib import Path
from typing import Any

from .diagram_engine import DiagramEngine, DiagramSpec
from .presentation_builder import SlideDeckSpec
from .schemas import ArtifactSpec
from .spreadsheet_builder import SpreadsheetSpec

DOCX_AVAILABLE = importlib.util.find_spec("docx") is not None
OPENPYXL_AVAILABLE = importlib.util.find_spec("openpyxl") is not None
PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None


class ArtifactUnavailableError(Exception):
    """Raised when an artifact type is requested but its optional
    rendering dependency is not installed. Never faked as success."""


def render_markdown(spec: ArtifactSpec) -> str:
    lines: list[str] = [f"# {spec.title}", ""]
    if spec.purpose:
        lines.append(f"_Purpose: {spec.purpose}_")
    if spec.audience:
        lines.append(f"_Audience: {spec.audience}_")
    lines.append("")
    for section in spec.content_sections:
        lines.append(f"## {section.heading}")
        if section.body:
            lines.append(section.body)
        for bullet in section.bullets:
            lines.append(f"- {bullet}")
        if section.table:
            headers = section.table.get("headers", [])
            rows = section.table.get("rows", [])
            if headers:
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in rows:
                    lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
        lines.append("")
    if spec.data_sources:
        lines.append("## Sources")
        for source in spec.data_sources:
            lines.append(f"- {source}")
    return "\n".join(lines)


class ArtifactRenderer:
    def __init__(self, output_root: Path):
        self.output_root = Path(output_root)

    def artifact_dir(self, artifact_id: str, version: str) -> Path:
        directory = self.output_root / artifact_id / version
        directory.mkdir(parents=True, exist_ok=True)
        return directory

    def render_markdown_file(self, spec: ArtifactSpec, version: str) -> Path:
        directory = self.artifact_dir(spec.id, version)
        path = directory / "report.md"
        path.write_text(render_markdown(spec), encoding="utf-8")
        return path

    def render_json_file(self, spec: ArtifactSpec, version: str, data: Any) -> Path:
        directory = self.artifact_dir(spec.id, version)
        path = directory / "data.json"
        path.write_text(json.dumps(data, indent=2, default=str), encoding="utf-8")
        return path

    def render_csv_file(self, spec: ArtifactSpec, version: str, headers: list[str], rows: list[list[Any]]) -> Path:
        directory = self.artifact_dir(spec.id, version)
        path = directory / "table.csv"
        with path.open("w", newline="", encoding="utf-8") as handle:
            writer = csv.writer(handle)
            writer.writerow(headers)
            writer.writerows(rows)
        return path

    def render_diagram_file(self, spec: ArtifactSpec, version: str, diagram_spec: DiagramSpec) -> Path:
        directory = self.artifact_dir(spec.id, version)
        mermaid_text = DiagramEngine().render_mermaid(diagram_spec)
        path = directory / "diagram.mmd"
        path.write_text(mermaid_text, encoding="utf-8")
        return path

    def render_docx_file(self, spec: ArtifactSpec, version: str) -> Path:
        if not DOCX_AVAILABLE:
            raise ArtifactUnavailableError("python-docx is not installed; DOCX rendering is unavailable")
        # Office renderers are intentionally imported only when their
        # capability is used.  Importing all three during Brain boot added
        # seconds to every desktop launch even when the user only wanted a
        # voice or memory command.
        from docx import Document

        directory = self.artifact_dir(spec.id, version)
        document = Document()
        document.add_heading(spec.title, level=1)
        if spec.purpose:
            document.add_paragraph(f"Purpose: {spec.purpose}")
        for section in spec.content_sections:
            document.add_heading(section.heading, level=2)
            if section.body:
                document.add_paragraph(section.body)
            for bullet in section.bullets:
                document.add_paragraph(bullet, style="List Bullet")
        path = directory / "report.docx"
        document.save(str(path))
        return path

    def render_spreadsheet_file(self, spec: ArtifactSpec, version: str, spreadsheet: SpreadsheetSpec) -> Path:
        if not OPENPYXL_AVAILABLE:
            raise ArtifactUnavailableError("openpyxl is not installed; spreadsheet rendering is unavailable")
        from openpyxl import Workbook

        directory = self.artifact_dir(spec.id, version)
        workbook = Workbook()
        workbook.remove(workbook.active)
        for sheet_spec in spreadsheet.sheets:
            sheet = workbook.create_sheet(sheet_spec.name[:31])
            sheet.append(sheet_spec.headers)
            for row in sheet_spec.rows:
                sheet.append(row)
            for cell_ref, formula in sheet_spec.formulas.items():
                sheet[cell_ref] = formula
            if sheet_spec.autofilter and sheet_spec.headers:
                last_col_index = len(sheet_spec.headers) - 1
                last_col = chr(ord("A") + last_col_index) if last_col_index < 26 else "Z"
                sheet.auto_filter.ref = f"A1:{last_col}{len(sheet_spec.rows) + 1}"
        path = directory / "workbook.xlsx"
        workbook.save(str(path))
        return path

    def render_presentation_file(self, spec: ArtifactSpec, version: str, deck: SlideDeckSpec) -> Path:
        if not PPTX_AVAILABLE:
            raise ArtifactUnavailableError("python-pptx is not installed; presentation rendering is unavailable")
        from pptx import Presentation

        directory = self.artifact_dir(spec.id, version)
        presentation = Presentation()
        title_layout = presentation.slide_layouts[0]
        bullet_layout = presentation.slide_layouts[1]
        title_slide = presentation.slides.add_slide(title_layout)
        title_slide.shapes.title.text = deck.title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"{deck.audience} — {deck.purpose}"
        for slide_spec in deck.slides:
            slide = presentation.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_spec.title
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, bullet in enumerate(slide_spec.bullets):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = bullet
        path = directory / "presentation.pptx"
        presentation.save(str(path))
        return path
